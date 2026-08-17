"""Generate the PME Care presentation as a .pptx.

    python -m pip install --target /tmp/pptxlibs python-pptx
    PYTHONPATH=/tmp/pptxlibs python docs/build_deck.py

Kept in the repository so the deck can be regenerated when the numbers change,
rather than being a binary somebody edits by hand and nobody can reproduce. The
figures at the top are the only things that need touching.

Screenshots come from docs/screenshots/. Regenerate those by running the app and
driving it in headless Chromium — see the interface notes in README.md.

Deliberately built with python-pptx rather than exported from the HTML deck: an
exported slide is a picture of a slide, which cannot be edited, re-worded or
re-ordered by whoever presents it.
"""

import pathlib
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = pathlib.Path(__file__).resolve().parent
SHOTS = HERE / "screenshots"
OUT = HERE.parent / "PME-Care.pptx"

# --- the figures the deck quotes -------------------------------------------
# Verified against the repository at build time; update together.
FACTS = {"roles": 4, "endpoints": 23, "tests": 118, "migrations": 5}

# --- palette, lifted from the application's own stylesheet ------------------
INK = RGBColor(0x17, 0x1B, 0x1E)
INK_SOFT = RGBColor(0x55, 0x63, 0x6B)
INK_FAINT = RGBColor(0x8A, 0x97, 0x9E)
GROUND = RGBColor(0xF2, 0xF5, 0xF4)
PANEL = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD6, 0xDC, 0xDA)
ACCENT = RGBColor(0x0B, 0x6B, 0x5F)
FIT = RGBColor(0x12, 0x68, 0x3A)
DUE = RGBColor(0x8A, 0x5A, 0x00)
OVERDUE = RGBColor(0xA0, 0x24, 0x18)

# Font names travel, font files do not: pick faces that ship with Office on both
# Windows and macOS so the deck does not silently re-flow on the presenter's
# laptop. Consolas for the small uppercase labels, matching the app's mono.
SANS = "Calibri"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)          # 16:9
MARGIN = Inches(0.9)
BODY_W = W - MARGIN * 2


def blank(prs: Presentation):
    """Add a slide with no placeholders and paint its background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = GROUND
    return slide


def textbox(slide, left, top, width, height):
    """A word-wrapping text box with no internal padding."""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    return frame


def write(frame, text, *, size, bold=False, color=INK, font=SANS,
          spacing=1.0, space_after=0, caps=False, align=PP_ALIGN.LEFT, first=False):
    """Append (or set) a paragraph of styled text."""
    para = frame.paragraphs[0] if first else frame.add_paragraph()
    para.alignment = align
    para.line_spacing = spacing
    para.space_after = Pt(space_after)
    run = para.add_run()
    run.text = text.upper() if caps else text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return para


def eyebrow(slide, text):
    """The small uppercase label above every headline."""
    frame = textbox(slide, MARGIN, Inches(0.62), BODY_W, Inches(0.3))
    write(frame, text, size=10.5, bold=True, color=ACCENT, font=MONO, caps=True, first=True)


def headline(slide, text, *, size=32):
    frame = textbox(slide, MARGIN, Inches(1.02), BODY_W, Inches(1.0))
    write(frame, text, size=size, bold=True, color=INK, spacing=0.95, first=True)


def rule(slide, top, width=None):
    """A hairline divider."""
    line = slide.shapes.add_shape(1, MARGIN, top, width or BODY_W, Emu(9525))
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    line.shadow.inherit = False
    return line


def card(slide, left, top, width, height, *, title, body, label=None, stripe=None):
    """A bordered panel: optional label, a title, and explanatory text."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = RULE
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    shape.text_frame.text = ""

    if stripe is not None:
        bar = slide.shapes.add_shape(1, left, top, Inches(0.045), height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = stripe
        bar.line.fill.background()
        bar.shadow.inherit = False

    pad = Inches(0.26)
    frame = textbox(slide, left + pad, top + pad, width - pad * 2, height - pad * 2)
    first = True
    if label:
        write(frame, label, size=8.5, bold=True, color=INK_FAINT, font=MONO,
              caps=True, space_after=5, first=True)
        first = False
    write(frame, title, size=14, bold=True, color=INK, space_after=4, first=first)
    write(frame, body, size=11, color=INK_SOFT, spacing=1.15)
    return shape


def stat(slide, left, top, width, value, label):
    """A figure with a rule down its left edge."""
    bar = slide.shapes.add_shape(1, left, top, Inches(0.04), Inches(0.85))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    frame = textbox(slide, left + Inches(0.18), top - Inches(0.04), width, Inches(0.9))
    write(frame, str(value), size=30, bold=True, color=INK, space_after=0, first=True)
    write(frame, label, size=9, bold=True, color=INK_SOFT, font=MONO, caps=True)


def picture(slide, name, *, top, height=None, width=None, left=None):
    """Place a screenshot, centred horizontally unless told otherwise."""
    path = SHOTS / f"{name}.png"
    if not path.exists():
        raise SystemExit(f"missing screenshot: {path}")
    kwargs = {}
    if height:
        kwargs["height"] = height
    if width:
        kwargs["width"] = width
    pic = slide.shapes.add_picture(str(path), left or Inches(0), top, **kwargs)
    if left is None:
        pic.left = int((W - pic.width) / 2)
    pic.line.color.rgb = RULE
    pic.line.width = Pt(0.75)
    return pic


def caption(slide, text, top):
    frame = textbox(slide, MARGIN, top, BODY_W, Inches(0.7))
    write(frame, text, size=11, color=INK_SOFT, spacing=1.2, first=True)


def bullets(slide, left, top, width, items, *, size=12):
    """A list where each item is "· text"; python-pptx has no real bullets."""
    frame = textbox(slide, left, top, width, Inches(3.4))
    for i, item in enumerate(items):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        para.line_spacing = 1.2
        para.space_after = Pt(9)
        dot = para.add_run()
        dot.text = "—  "
        dot.font.size = Pt(size)
        dot.font.color.rgb = ACCENT
        dot.font.name = SANS
        run = para.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = INK_SOFT
        run.font.name = SANS


def build() -> pathlib.Path:
    """Write the twelve slides."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 01 — title -------------------------------------------------------------
    s = blank(prs)
    eyebrow(s, "Occupational health · internal software")
    frame = textbox(s, MARGIN, Inches(2.15), BODY_W, Inches(1.3))
    write(frame, "PME Care", size=64, bold=True, color=INK, spacing=0.9, first=True)
    frame = textbox(s, MARGIN, Inches(3.5), Inches(7.2), Inches(1.2))
    write(frame, "Periodic medical examination tracking for an industrial workforce "
                 "— who is due, who examined them, and what was decided.",
          size=17, color=INK_SOFT, spacing=1.25, first=True)
    labels = [("roles", "Roles"), ("endpoints", "Endpoints"),
              ("tests", "Tests passing"), ("migrations", "Schema versions")]
    for i, (key, label) in enumerate(labels):
        stat(s, MARGIN + Inches(2.95) * i, Inches(5.25), Inches(2.6), FACTS[key], label)

    # 02 — the problem -------------------------------------------------------
    s = blank(prs)
    eyebrow(s, "The problem")
    headline(s, "A lapsed examination is silent")
    frame = textbox(s, MARGIN, Inches(2.3), Inches(5.9), Inches(3))
    write(frame, "Periodic medical examinations are periodic by law. The failure mode is "
                 "not dramatic — nobody refuses an examination. Someone simply falls off "
                 "the schedule, and nothing says so.",
          size=13.5, color=INK_SOFT, spacing=1.3, space_after=12, first=True)
    write(frame, "A register of past examinations cannot answer the only question that "
                 "matters: who is overdue right now? Answering it from paper means reading "
                 "every record and doing arithmetic on every date.",
          size=13.5, color=INK_SOFT, spacing=1.3)
    right = MARGIN + Inches(6.5)
    card(s, right, Inches(2.3), Inches(5.05), Inches(1.5), label="Before",
         title="A record of the past",
         body="Examinations are filed after the fact. Compliance is asserted at audit "
              "time and reconstructed by hand.")
    card(s, right, Inches(4.0), Inches(5.05), Inches(1.6), label="After",
         title="A queue for today", stripe=ACCENT,
         body="The system computes when each next examination falls due and puts the "
              "people behind that date on one screen, worst first.")

    # 03 — roles -------------------------------------------------------------
    s = blank(prs)
    eyebrow(s, "Who uses it")
    headline(s, "Four roles, four different screens")
    roles = [
        ("Health Team", "Keeps the schedule",
         "Registers employees, sees who has lapsed, books examinations, cancels with a "
         "reason. Cannot record a clinical decision."),
        ("Doctor", "Makes the decision",
         "Works a list of booked examinations and records the outcome against the "
         "person's history. Cannot register or cancel."),
        ("Employee", "Sees only themselves",
         "One screen: am I fit, when is my next examination, what was said. No lists, "
         "and no way to reach another person's record."),
        ("Administrator", "Accounts and accountability",
         "Creates and revokes staff accounts, and reads the audit trail. Cannot lock "
         "themselves out."),
    ]
    for i, (label, title, body) in enumerate(roles):
        col, row = i % 2, i // 2
        card(s, MARGIN + (BODY_W / 2 + Inches(0.15)) * col,
             Inches(2.25) + Inches(1.95) * row,
             BODY_W / 2 - Inches(0.15), Inches(1.75),
             label=label, title=title, body=body)
    frame = textbox(s, MARGIN, Inches(6.3), BODY_W, Inches(0.5))
    write(frame, "Every role's permissions are enforced by the server on each request. "
                 "The interface simply stops offering what would be refused.",
          size=10.5, color=INK_FAINT, first=True)

    # 04 — the cycle ---------------------------------------------------------
    s = blank(prs)
    eyebrow(s, "The cycle")
    headline(s, "One loop, and it closes itself")
    steps = [("Register", "Health Team", None), ("Falls due", "computed", OVERDUE),
             ("Book", "Health Team", None), ("Examine", "Doctor", None),
             ("Record outcome", "Fit / unfit", FIT), ("Next due date", "automatic", DUE)]
    box_w = (BODY_W - Inches(0.5)) / 6
    for i, (title, sub, stripe) in enumerate(steps):
        left = MARGIN + (box_w + Inches(0.1)) * i
        shape = s.shapes.add_shape(1, left, Inches(2.35), box_w, Inches(1.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL
        shape.line.color.rgb = RULE
        shape.line.width = Pt(0.75)
        shape.shadow.inherit = False
        if stripe is not None:
            bar = s.shapes.add_shape(1, left, Inches(2.35), Inches(0.04), Inches(1.05))
            bar.fill.solid()
            bar.fill.fore_color.rgb = stripe
            bar.line.fill.background()
            bar.shadow.inherit = False
        frame = textbox(s, left + Inches(0.16), Inches(2.55), box_w - Inches(0.3), Inches(0.7))
        write(frame, title, size=12, bold=True, color=INK, space_after=2, first=True)
        write(frame, sub, size=8.5, color=INK_FAINT, font=MONO)
    frame = textbox(s, MARGIN, Inches(3.95), Inches(5.9), Inches(2))
    write(frame, "Completing an examination is what starts the next cycle. A Fit outcome "
                 "is good for a year and the employee leaves the queue; a provisional one "
                 "returns them within months.",
          size=13.5, color=INK_SOFT, spacing=1.3, first=True)
    frame = textbox(s, MARGIN + Inches(6.5), Inches(3.95), Inches(5.05), Inches(2))
    write(frame, "An Unfit outcome sets no due date at all. That is a case to manage, not "
                 "a booking to make — inventing a routine recall would quietly downgrade a "
                 "serious finding.",
          size=13.5, color=INK_SOFT, spacing=1.3, first=True)

    # 05–06 — the two landscape screens, in role order ----------------------
    for eb, hl, shot, cap in [
        ("Health Team · the landing screen", "Who needs an examination booked",
         "compliance", "The first screen is the work, not a search box. Counts across the "
         "top, the people behind them below, and Book today on every row — booking someone "
         "removes them from the list, so what remains is what is left to do."),
        ("Doctor · recording an outcome", "The decision, beside the trend",
         "examine", "Previous readings sit beside the form, not behind a click — a fitness "
         "decision is a judgement about a trend. Remarks become mandatory the moment the "
         "decision is anything other than Fit."),
    ]:
        s = blank(prs)
        eyebrow(s, eb)
        headline(s, hl, size=28)
        picture(s, shot, top=Inches(1.95), height=Inches(4.0))
        caption(s, cap, Inches(6.15))

    # 07 — employee, portrait, so it gets its own treatment ------------------
    s = blank(prs)
    eyebrow(s, "Employee · the whole interface")
    headline(s, "One question, answered", size=28)
    picture(s, "employee", top=Inches(1.95), height=Inches(4.6), left=MARGIN)
    frame = textbox(s, MARGIN + Inches(4.4), Inches(2.1), Inches(7.6), Inches(4))
    write(frame, "Am I cleared for work? The answer is the first thing on the screen, in "
                 "words, before any detail. Below it: their own particulars and their own "
                 "history.",
          size=14, color=INK_SOFT, spacing=1.3, space_after=14, first=True)
    write(frame, "Deliberately absent: any list, any search, any other person, and any "
                 "identifier they would have to type. An employee is shown their record "
                 "because the system already knows whose it is.",
          size=14, color=INK_SOFT, spacing=1.3, space_after=14)
    write(frame, '"Not yet examined" is worded distinctly from "Fit" — the two mean very '
                 "different things to someone about to start a shift.",
          size=11, color=INK_FAINT, spacing=1.25)

    # 08 — the audit trail, last of the four screens -------------------------
    s = blank(prs)
    eyebrow(s, "Administrator · accountability")
    headline(s, "Every change, and no clinical values", size=28)
    picture(s, "audit", top=Inches(1.95), height=Inches(4.0))
    caption(s, "The trail records which fields changed, never what they changed to. A "
               "blood-pressure reading in an audit table would be medical data sitting "
               "outside the controls that protect the examination itself.", Inches(6.15))

    # 09 — the rules ---------------------------------------------------------
    s = blank(prs)
    eyebrow(s, "Rules the software enforces")
    headline(s, "Refusals are the feature")
    rows = [
        ("Book a second open examination for one person", "Refused",
         "Two open examinations make “their current status” ambiguous. Enforced by "
         "the database, not a check two clicks could race past."),
        ("Record Unfit with no remarks", "Refused",
         "An outcome with consequences has to carry the reason for it."),
        ("Cancel an examination without a reason", "Refused",
         "A cancellation is a gap in the record that someone must be able to explain."),
        ("Employee opens another employee's record", "Not found",
         "Not “forbidden” — that would confirm the record exists."),
        ("Administrator deactivates their own account", "Refused",
         "That is how a deployment ends up with nobody able to administer it."),
        ("Re-open a completed examination", "Refused",
         "What was decided, and when, stays fixed. A new examination is booked instead."),
    ]
    widths = [Inches(4.2), Inches(1.5), Inches(5.85)]
    top = Inches(2.15)
    for i, head in enumerate(("Attempt", "Answer", "Why")):
        frame = textbox(s, MARGIN + sum(widths[:i]), top, widths[i], Inches(0.3))
        write(frame, head, size=8.5, bold=True, color=INK_FAINT, font=MONO, caps=True, first=True)
    y = top + Inches(0.34)
    for attempt, answer, why in rows:
        rule(s, y)
        for i, (text, size, color, bold) in enumerate((
                (attempt, 11.5, INK, False), (answer, 11.5, OVERDUE, True), (why, 11, INK_SOFT, False))):
            frame = textbox(s, MARGIN + sum(widths[:i]) , y + Inches(0.12),
                            widths[i] - Inches(0.25), Inches(0.6))
            write(frame, text, size=size, bold=bold, color=color, spacing=1.15, first=True)
        y += Inches(0.72)
    rule(s, y)

    # 10 — how it is built ---------------------------------------------------
    s = blank(prs)
    eyebrow(s, "How it is built")
    headline(s, "Boring on purpose")
    bullets(s, MARGIN, Inches(2.3), Inches(5.9), [
        "Python and PostgreSQL — chosen for longevity, not novelty.",
        "No build step in the interface. Deploying it is copying files; there is no "
        "toolchain to break in two years.",
        "One process serves the interface and the data together, on one port.",
        "Runs entirely on your own hardware. No data leaves the site.",
    ], size=13)
    right = MARGIN + Inches(6.5)
    card(s, right, Inches(2.3), Inches(5.05), Inches(1.75), title="Schema changes are versioned",
         body="Every change is a numbered migration, tested to apply and to reverse. A "
              "change that cannot be undone is a deployment that cannot be rolled back.")
    card(s, right, Inches(4.25), Inches(5.05), Inches(1.6), title="History outlives the record",
         body="Nothing is hard-deleted. Retiring an employee keeps their examinations and "
              "their audit trail intact.")

    # 11 — status ------------------------------------------------------------
    s = blank(prs)
    eyebrow(s, "Where it stands")
    headline(s, "Verified, and honestly incomplete")
    frame = textbox(s, MARGIN, Inches(2.2), Inches(5.9), Inches(0.3))
    write(frame, "Working and tested", size=9.5, bold=True, color=FIT, font=MONO, caps=True, first=True)
    bullets(s, MARGIN, Inches(2.6), Inches(5.9), [
        "The full cycle, end to end, against a real database.",
        f"{FACTS['tests']} automated tests, run on every change, including the database "
        "migrations in both directions.",
        "Every screen driven in a real browser as all four roles.",
        "Accounts, roles, revocation, and the audit trail.",
    ])
    right = MARGIN + Inches(6.5)
    frame = textbox(s, right, Inches(2.2), Inches(5.05), Inches(0.3))
    write(frame, "Not built yet", size=9.5, bold=True, color=DUE, font=MONO, caps=True, first=True)
    bullets(s, right, Inches(2.6), Inches(5.05), [
        "Sign-in from a browser. The server verifies credentials correctly; the login "
        "page that obtains them is not finished. This is the one thing between here and "
        "real users.",
        "Reminders. Due dates are computed and visible, but nothing sends anything to "
        "anybody yet.",
        "Attachments, and a richer set of recorded measurements.",
    ])

    # 12 — decisions ---------------------------------------------------------
    s = blank(prs)
    eyebrow(s, "What I need from you")
    headline(s, "Three decisions, then it ships")
    decisions = [
        ("Decision one · highest priority", "Are the examination intervals correct?",
         "The software recalls a fit employee after 12 months and a provisionally fit one "
         "after 3. I chose those numbers. For a hazardous process they may be set by "
         "regulation — and they should be right before real records exist. They are "
         "configuration, not code.", OVERDUE),
        ("Decision two", "How do people sign in?",
         "Sign-in currently depends on an external service, so an internet connection is "
         "needed to log in to an otherwise self-contained system. Using the existing plant "
         "directory instead would remove that, and mean no new passwords for anyone.", DUE),
        ("Decision three", "Who owns backups?",
         "On your own hardware this is nobody's job until it is named as somebody's. A "
         "nightly copy is scripted; a restore has never been performed. An untested backup "
         "is not a backup.", ACCENT),
    ]
    col_w = (BODY_W - Inches(0.4)) / 3
    for i, (label, title, body, stripe) in enumerate(decisions):
        card(s, MARGIN + (col_w + Inches(0.2)) * i, Inches(2.25), col_w, Inches(3.3),
             label=label, title=title, body=body, stripe=stripe)
    frame = textbox(s, MARGIN, Inches(6.0), BODY_W, Inches(0.4))
    write(frame, "Source, documentation and deployment runbook: "
                 "github.com/nichepah/pme-care",
          size=10.5, color=INK_FAINT, font=MONO, first=True)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    try:
        path = build()
    except ImportError:
        sys.exit("python-pptx is not installed — see this file's docstring.")
    print(f"wrote {path} ({path.stat().st_size / 1024:.0f} KB)")
